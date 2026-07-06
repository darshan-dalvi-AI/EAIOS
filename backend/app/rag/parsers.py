"""Multi-format document parsing → list of text blocks with section/page metadata.
Layout-aware upgrades (unstructured / Docling / PaddleOCR) plug in behind the same interface."""
import csv
import os

Block = dict  # {"section": str, "page": int, "text": str}


def parse_file(path: str, doc_type: str) -> list[Block]:
    parser = {
        "pdf": _parse_pdf,
        "docx": _parse_docx,
        "pptx": _parse_pptx,
        "xlsx": _parse_xlsx,
        "csv": _parse_csv,
        "txt": _parse_txt,
        "image": _parse_image,
    }.get(doc_type)
    if parser is None:
        raise ValueError(f"No parser for doc_type={doc_type}")
    blocks = [b for b in parser(path) if b["text"].strip()]
    if not blocks:
        blocks = [{"section": "", "page": 0, "text": f"[No extractable text in {os.path.basename(path)}]"}]
    return blocks


def _parse_pdf(path: str) -> list[Block]:
    from pypdf import PdfReader

    reader = PdfReader(path)
    return [
        {"section": "", "page": i + 1, "text": (page.extract_text() or "")}
        for i, page in enumerate(reader.pages)
    ]


def _parse_docx(path: str) -> list[Block]:
    import docx

    blocks: list[Block] = []
    section = ""
    buffer: list[str] = []
    for para in docx.Document(path).paragraphs:
        text = para.text.strip()
        if not text:
            continue
        if para.style.name.startswith("Heading"):
            if buffer:
                blocks.append({"section": section, "page": 0, "text": "\n".join(buffer)})
                buffer = []
            section = text
        else:
            buffer.append(text)
    if buffer:
        blocks.append({"section": section, "page": 0, "text": "\n".join(buffer)})
    return blocks


def _parse_pptx(path: str) -> list[Block]:
    from pptx import Presentation

    blocks: list[Block] = []
    for i, slide in enumerate(Presentation(path).slides, start=1):
        texts = [
            shape.text_frame.text.strip()
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        if texts:
            blocks.append({"section": texts[0][:80], "page": i, "text": "\n".join(texts)})
    return blocks


def _parse_xlsx(path: str, max_rows: int = 200) -> list[Block]:
    from openpyxl import load_workbook

    wb = load_workbook(path, read_only=True, data_only=True)
    blocks: list[Block] = []
    for sheet in wb.worksheets:
        lines: list[str] = []
        for r, row in enumerate(sheet.iter_rows(values_only=True)):
            if r >= max_rows:
                lines.append(f"... ({sheet.max_row - max_rows} more rows)")
                break
            cells = [str(c) for c in row if c is not None]
            if cells:
                lines.append(" | ".join(cells))
        if lines:
            blocks.append({"section": sheet.title, "page": 0, "text": "\n".join(lines)})
    return blocks


def _parse_csv(path: str, max_rows: int = 300) -> list[Block]:
    lines: list[str] = []
    with open(path, newline="", encoding="utf-8", errors="replace") as f:
        for r, row in enumerate(csv.reader(f)):
            if r >= max_rows:
                lines.append("... (truncated)")
                break
            lines.append(" | ".join(row))
    return [{"section": os.path.basename(path), "page": 0, "text": "\n".join(lines)}]


def _parse_txt(path: str) -> list[Block]:
    with open(path, encoding="utf-8", errors="replace") as f:
        content = f.read()
    blocks: list[Block] = []
    section = ""
    buffer: list[str] = []
    for line in content.splitlines():
        if line.startswith("#"):  # markdown heading → section boundary
            if buffer:
                blocks.append({"section": section, "page": 0, "text": "\n".join(buffer)})
                buffer = []
            section = line.lstrip("# ").strip()
        else:
            buffer.append(line)
    if buffer:
        blocks.append({"section": section, "page": 0, "text": "\n".join(buffer)})
    return blocks


def _parse_image(path: str) -> list[Block]:
    """Vision pipeline: OCR (pytesseract) + VLM caption (local Ollama vision model), whichever is available."""
    blocks: list[Block] = []
    try:
        import pytesseract
        from PIL import Image

        text = pytesseract.image_to_string(Image.open(path)).strip()
        if text:
            blocks.append({"section": "OCR", "page": 0, "text": text})
    except ImportError:
        pass

    from app.llm.provider import caption_image

    caption = caption_image(path)
    if caption:
        blocks.append({"section": "AI Caption", "page": 0, "text": caption})

    if not blocks:
        blocks.append({
            "section": "",
            "page": 0,
            "text": f"[Image '{os.path.basename(path)}' uploaded — install pytesseract for OCR or pull a vision model (e.g. `ollama pull llava`) for AI captioning]",
        })
    return blocks
