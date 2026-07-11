"""Advanced document parsing — structured table extraction → SQL Agent.

Complex/nested tables in uploaded files are extracted as *structured data*
and materialized into REAL SQL tables (``dt_<doc>_<n>``) in the platform
database, bypassing text chunking entirely. The SQL Agent sees them in its
schema and can aggregate/filter them precisely — far more accurate than
retrieving a table as flattened prose.

Extraction per format:
- **docx / pptx** — native table objects (nested docx tables are walked
  recursively, depth-first).
- **xlsx / csv**  — sheets/rows directly.
- **pdf / txt**   — line-grid heuristics on extracted text: markdown pipes,
  tabs, or 2+-space column alignment across ≥3 consecutive lines.

Each materialized table also emits one small summary block (schema + sample
rows) that IS chunked and embedded, so hybrid RAG can cite the table and the
LLM learns the physical table name to hand to the SQL Agent.
"""
import csv
import json
import logging
import re
from dataclasses import dataclass, field

from sqlalchemy import text as sqltext
from sqlalchemy.orm import Session

from app.models import DataTable, Document

log = logging.getLogger("eaios.tables")

MAX_TABLES = 10
MAX_ROWS = 500
MAX_COLS = 30
MIN_DATA_ROWS = 2   # header + ≥2 data rows to bother materializing
MIN_COLS = 2

Block = dict  # {"section": str, "page": int, "text": str}


@dataclass
class RawTable:
    title: str = ""
    columns: list[str] = field(default_factory=list)
    rows: list[list[str]] = field(default_factory=list)


# ── identifier & type hygiene ────────────────────────────────────────────
def sanitize_ident(name: str, used: set[str], fallback: str) -> str:
    ident = re.sub(r"[^a-z0-9_]+", "_", str(name).strip().lower()).strip("_")[:40]
    if not ident or not re.match(r"^[a-z_]", ident):
        ident = fallback
    base, n = ident, 2
    while ident in used:
        ident, n = f"{base}_{n}", n + 1
    used.add(ident)
    return ident


_NUM_RX = re.compile(r"^-?[\d,]+(\.\d+)?%?$")


def _numeric(value: str) -> float | None:
    v = str(value).strip().replace(",", "").rstrip("%").lstrip("$₹€£").strip()
    if not v or not _NUM_RX.match(str(value).strip().lstrip("$₹€£").strip()):
        return None
    try:
        return float(v)
    except ValueError:
        return None


def infer_types(rows: list[list[str]], ncols: int) -> list[str]:
    """Per-column SQL type: INTEGER / REAL / TEXT (valid in SQLite AND Postgres)."""
    types: list[str] = []
    for c in range(ncols):
        values = [str(r[c]).strip() for r in rows if c < len(r) and str(r[c]).strip()]
        if not values:
            types.append("TEXT")
            continue
        nums = [_numeric(v) for v in values]
        hits = [n for n in nums if n is not None]
        if len(hits) >= max(1, int(0.9 * len(values))):
            types.append("INTEGER" if all(float(h).is_integer() for h in hits) else "REAL")
        else:
            types.append("TEXT")
    return types


def _coerce(value: str, sql_type: str):
    v = str(value).strip()
    if sql_type == "TEXT":
        return v
    n = _numeric(v)
    if n is None:
        return None
    return int(n) if sql_type == "INTEGER" else n


# ── text-grid detection (pdf / txt / markdown) ───────────────────────────
_MD_SEP = re.compile(r"^\s*\|?[\s:|-]+\|[\s:|-]+\|?\s*$")


def _split_line(line: str) -> tuple[str, list[str]]:
    """→ (delimiter_kind, cells). kind ∈ pipe|tab|space|'' (not a row)."""
    if line.count("|") >= 2:
        cells = [c.strip() for c in line.strip().strip("|").split("|")]
        if len(cells) >= MIN_COLS:
            return "pipe", cells
    if line.count("\t") >= 1:
        cells = [c.strip() for c in line.split("\t")]
        if len([c for c in cells if c]) >= MIN_COLS:
            return "tab", cells
    cells = [c.strip() for c in re.split(r"\s{2,}", line.strip()) if c.strip()]
    if len(cells) >= 3:  # stricter for space-aligned — prose has 2-col false positives
        return "space", cells
    return "", []


def detect_text_tables(text: str) -> list[RawTable]:
    """Find runs of ≥3 consecutive same-delimiter lines that look like a grid.
    Short rows are padded (handles merged/nested header cells)."""
    tables: list[RawTable] = []
    run: list[list[str]] = []
    run_kind = ""

    def flush() -> None:
        nonlocal run, run_kind
        if len(run) >= MIN_DATA_ROWS + 1:
            width = max(len(r) for r in run)
            if MIN_COLS <= width:
                rows = [r + [""] * (width - len(r)) for r in run]
                tables.append(RawTable(columns=rows[0][:MAX_COLS],
                                       rows=[r[:MAX_COLS] for r in rows[1:MAX_ROWS + 1]]))
        run, run_kind = [], ""

    for line in text.splitlines():
        if _MD_SEP.match(line):  # markdown |---|---| separator: keep run alive
            if run_kind == "pipe":
                continue
        kind, cells = _split_line(line)
        if kind and (not run_kind or kind == run_kind):
            run_kind = run_kind or kind
            run.append(cells)
        else:
            flush()
            if kind:
                run_kind = kind
                run.append(cells)
    flush()
    return tables


# ── per-format extraction ────────────────────────────────────────────────
def extract_tables(path: str, doc_type: str) -> list[RawTable]:
    try:
        extractor = {
            "pdf": _from_pdf, "docx": _from_docx, "pptx": _from_pptx,
            "xlsx": _from_xlsx, "csv": _from_csv, "txt": _from_txt,
        }.get(doc_type)
        if extractor is None:
            return []
        found = extractor(path)[:MAX_TABLES]
        return [t for t in found if len(t.rows) >= MIN_DATA_ROWS and len(t.columns) >= MIN_COLS]
    except Exception:  # noqa: BLE001 — table extraction is additive, never blocks ingest
        log.exception("table extraction failed for %s", path)
        return []


def _from_pdf(path: str) -> list[RawTable]:
    from pypdf import PdfReader

    tables: list[RawTable] = []
    for i, page in enumerate(PdfReader(path).pages, start=1):
        for n, t in enumerate(detect_text_tables(page.extract_text() or ""), start=1):
            t.title = f"p.{i} table {n}"
            tables.append(t)
    return tables


def _walk_docx_tables(container, out: list[RawTable], depth: int = 0) -> None:
    """Depth-first over docx tables — tables nested inside cells become their
    own structured tables (the 'complex, nested tables' case)."""
    for ti, table in enumerate(container.tables, start=1):
        grid = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        if grid and len(grid[0]) >= MIN_COLS:
            title = f"table {ti}" + (" (nested)" if depth else "")
            out.append(RawTable(title=title, columns=grid[0][:MAX_COLS],
                                rows=[r[:MAX_COLS] for r in grid[1:MAX_ROWS + 1]]))
        if depth < 3:
            for row in table.rows:
                for cell in row.cells:
                    if cell.tables:
                        _walk_docx_tables(cell, out, depth + 1)


def _from_docx(path: str) -> list[RawTable]:
    import docx

    out: list[RawTable] = []
    _walk_docx_tables(docx.Document(path), out)
    return out


def _from_pptx(path: str) -> list[RawTable]:
    from pptx import Presentation

    out: list[RawTable] = []
    for i, slide in enumerate(Presentation(path).slides, start=1):
        for shape in slide.shapes:
            if getattr(shape, "has_table", False) and shape.has_table:
                grid = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                if grid:
                    out.append(RawTable(title=f"slide {i} table", columns=grid[0][:MAX_COLS],
                                        rows=[r[:MAX_COLS] for r in grid[1:MAX_ROWS + 1]]))
    return out


def _from_xlsx(path: str) -> list[RawTable]:
    from openpyxl import load_workbook

    out: list[RawTable] = []
    for sheet in load_workbook(path, read_only=True, data_only=True).worksheets:
        grid = [["" if c is None else str(c) for c in row]
                for row in sheet.iter_rows(values_only=True)][:MAX_ROWS + 1]
        grid = [r for r in grid if any(c.strip() for c in r)]
        if grid and len(grid[0]) >= MIN_COLS:
            out.append(RawTable(title=sheet.title, columns=grid[0][:MAX_COLS],
                                rows=[r[:MAX_COLS] for r in grid[1:]]))
    return out


def _from_csv(path: str) -> list[RawTable]:
    with open(path, newline="", encoding="utf-8", errors="replace") as f:
        grid = [row for _, row in zip(range(MAX_ROWS + 1), csv.reader(f))]
    grid = [r for r in grid if any(str(c).strip() for c in r)]
    if not grid or len(grid[0]) < MIN_COLS:
        return []
    return [RawTable(columns=[str(c) for c in grid[0][:MAX_COLS]],
                     rows=[[str(c) for c in r[:MAX_COLS]] for r in grid[1:]])]


def _from_txt(path: str) -> list[RawTable]:
    with open(path, encoding="utf-8", errors="replace") as f:
        return detect_text_tables(f.read())


# ── materialization ──────────────────────────────────────────────────────
def drop_for_document(db: Session, doc_id: str) -> int:
    """Reindex/delete hygiene: drop this document's physical tables + metadata."""
    dropped = 0
    for dt in db.query(DataTable).filter(DataTable.document_id == doc_id).all():
        db.execute(sqltext(f'DROP TABLE IF EXISTS "{dt.table_name}"'))
        db.delete(dt)
        dropped += 1
    db.flush()
    return dropped


def materialize(db: Session, doc: Document, tables: list[RawTable]) -> list[Block]:
    """Create physical SQL tables + DataTable metadata; return one summary
    block per table for the RAG index."""
    drop_for_document(db, doc.id)
    blocks: list[Block] = []

    for i, raw in enumerate(tables, start=1):
        used: set[str] = set()
        cols = [sanitize_ident(c, used, f"col_{n + 1}") for n, c in enumerate(raw.columns)]
        types = infer_types(raw.rows, len(cols))
        table_name = f"dt_{doc.id[:8]}_{i}"

        ddl_cols = ", ".join(f'"{c}" {t}' for c, t in zip(cols, types))
        db.execute(sqltext(f'DROP TABLE IF EXISTS "{table_name}"'))
        db.execute(sqltext(f'CREATE TABLE "{table_name}" ({ddl_cols})'))

        params = [
            {f"c{n}": _coerce(row[n] if n < len(row) else "", types[n]) for n in range(len(cols))}
            for row in raw.rows
        ]
        placeholders = ", ".join(f":c{n}" for n in range(len(cols)))
        col_list = ", ".join(f'"{c}"' for c in cols)
        if params:
            db.execute(sqltext(f'INSERT INTO "{table_name}" ({col_list}) VALUES ({placeholders})'), params)

        db.add(DataTable(
            document_id=doc.id, doc_title=doc.title, table_name=table_name,
            title=raw.title or f"table {i}", source=doc.doc_type, row_count=len(raw.rows),
            columns=json.dumps([{"name": c, "type": t} for c, t in zip(cols, types)]),
        ))

        sample = "\n".join(" | ".join(str(v) for v in row[:8]) for row in raw.rows[:3])
        blocks.append({
            "section": f"Structured table: {raw.title or f'table {i}'}",
            "page": 0,
            "text": (
                f"[STRUCTURED TABLE {table_name} — '{raw.title or f'table {i}'}' from '{doc.title}'] "
                f"{len(raw.rows)} rows × {len(cols)} columns, materialized in the SQL database. "
                f"The SQL Agent can query it directly, e.g. SELECT * FROM {table_name} LIMIT 10\n"
                f"columns: {', '.join(f'{c} ({t})' for c, t in zip(cols, types))}\n"
                f"sample rows:\n{sample}"
            ),
        })

    db.flush()
    return blocks


def ingest_tables(db: Session, doc: Document, path: str) -> list[Block]:
    """Pipeline entry point: extract + materialize; returns summary blocks."""
    tables = extract_tables(path, doc.doc_type)
    if not tables:
        drop_for_document(db, doc.id)  # file may have lost its tables on re-upload
        return []
    blocks = materialize(db, doc, tables)
    log.info("Materialized %d structured table(s) for %s", len(blocks), doc.filename)
    return blocks
